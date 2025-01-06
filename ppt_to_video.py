import os
import re
from pptx import Presentation
from gtts import gTTS
from moviepy.editor import *
import comtypes.client
import comtypes
from pptx.enum.shapes import PP_MEDIA_TYPE, MSO_SHAPE_TYPE
import subprocess
import zipfile


# Define the folder paths at the top
TEMP_IMAGES_FOLDER = 'temp_images'
TEMP_AUDIO_FOLDER = 'temp_audio'
TEMP_VIDEOS_FOLDER = 'temp_videos'

def create_temp_folders():
    if not os.path.exists(TEMP_IMAGES_FOLDER):
        os.mkdir(TEMP_IMAGES_FOLDER)
    if not os.path.exists(TEMP_AUDIO_FOLDER):
        os.mkdir(TEMP_AUDIO_FOLDER)
    if not os.path.exists(TEMP_VIDEOS_FOLDER):
        os.mkdir(TEMP_VIDEOS_FOLDER)


def cleanup_temp_dirs():
    """Clean up temporary directories with error handling"""
    print("Cleaning up temporary directories...")
    for directory in [TEMP_IMAGES_FOLDER, TEMP_AUDIO_FOLDER, TEMP_VIDEOS_FOLDER]:
        if os.path.exists(directory):
            try:
                for file in os.listdir(directory):
                    try:
                        file_path = os.path.join(directory, file)
                        if os.path.isfile(file_path):
                            os.chmod(file_path, 0o777)  # Give full permissions
                            os.remove(file_path)
                    except Exception as e:
                        print(f"Warning: Could not remove file {file}: {e}")
                os.rmdir(directory)
            except Exception as e:
                print(f"Warning: Could not remove directory {directory}: {e}")

def get_slide_settings(notes):
    """
    Get the following settings from the slide notes (if available):
    min_time, pause_time_at_end, ai_voice
    Return the settings as a dictionary and the cleaned notes without the settings.
    """

    settings = {
        "min_time": None,
        "pause_time_at_end": None,
        "ai_voice": None
    }

    # check for min_time in notes
    match = re.search(r'\{\{min_time:(\d+)\}\}', notes)
    if match:
        min_time = int(match.group(1).strip())
        notes = re.sub(r'\{\{min_time:\d+\}\}', '', notes)
        settings["min_time"] = min_time

    # check for pause_time_at_end in notes
    match = re.search(r'\{\{pause_time_at_end:(\d+)\}\}', notes)
    if match:
        pause_time_at_end = int(match.group(1).strip())
        notes = re.sub(r'\{\{pause_time_at_end:\d+\}\}', '', notes)
        settings["pause_time_at_end"] = pause_time_at_end

    # check for ai_voice in notes
    match = re.search(r'\{\{ai_voice:(.*?)\}\}', notes)
    if match:
        ai_voice = match.group(1).strip()
        notes = re.sub(r'\{\{ai_voice:(.*?)\}\}', '', notes)
        settings["ai_voice"] = ai_voice

    return settings, notes

def remove_comments(notes=None):
    return re.sub(r'\{\*.*?\*\}', '', notes, flags=re.DOTALL) if notes is not None else None # DOTAALL flag is used to match multilines


def speed_up_audio_ffmpeg(input_path, output_path, speed_factor=1.25):
    """
    NOT recommended to speed up audio files individually before concatenating them. 
    First concatenate the audio/video files and then speed up the final file. Otherwise you might hear some glitches.
    """
    subprocess.run([
        "ffmpeg", 
        "-y", #overwrite file without asking
        "-i", input_path, 
        "-filter:a", f"atempo={speed_factor}",
        "-b:a", "320k",
        "-q:a", "1",
        "-vn",                # no video
        output_path
    ], check=True)


def speed_up_video_ffmpeg(input_path, output_path, speed_factor=1.25):
    command = [
        "ffmpeg",
        "-y", #overwrite file without asking
        "-i", input_path,
        "-filter_complex", f"[0:v]setpts=0.8*PTS[v];[0:a]atempo={speed_factor}[a]",
        "-map", "[v]",
        "-map", "[a]",
        output_path
    ]
    subprocess.run(command, check=True)
def extract_videos_from_slides(ppt_path):

    videos = []
    output_folder = TEMP_VIDEOS_FOLDER
    # Ensure the output folder exists
    os.makedirs(output_folder, exist_ok=True)

    # Open the .pptx file as a ZIP archive to access embedded media
    with zipfile.ZipFile(ppt_path, 'r') as pptx_zip:
        videos = [x for x in pptx_zip.namelist() if "mp4" in x]

        for video in videos:
            video_filename = os.path.basename(video)
            video_path = os.path.join(output_folder, video_filename)
            with open(video_path, "wb") as video_file:
                video_file.write(pptx_zip.read(video))
            print(f"Extracted video: {video_path}")

    return videos

def convert_ppt_to_video(openai_client, ppt_path, output_dir="output", output_video="output.mp4", provider="google", language='en', accent='com', openai_voice='alloy', extra_settings=None):
    # Default Settings
    min_time_per_slide = 6,
    pause_time_at_end = 1
    speed_factor = 1
    fps=30
    slides_numbers_to_process = []
    
    if extra_settings is not None:
        print("Extra settings: ", extra_settings)
        min_time_per_slide = extra_settings.get("min_time_per_slide", 6)
        pause_time_at_end = extra_settings.get("pause_time_at_end", 1)
        speed_factor = extra_settings.get("speed_factor", 1)
        fps = extra_settings.get("fps", 30)
        slide_numbers = extra_settings.get("slide_numbers", None)

        if slide_numbers:
            # get slide numbers to process. Support for example 1,2,3 for slides 1,2,3 and also 1-5 for all slides from 1 to 5 inclusive
            for part in slide_numbers.split(","):
                if "-" in part:
                    start, end = part.split("-")
                    slides_numbers_to_process.extend(range(int(start), int(end) + 1))
                else:
                    slides_numbers_to_process.append(int(part))
    
        print("Slides to process:", slides_numbers_to_process)
    try:
        clips = []
        # Create temp directories
        create_temp_folders()
        if not os.path.exists(output_dir):
            os.mkdir(output_dir)

        print("File path: ", ppt_path)
        if provider == "openai":
            print("Using OpenAI with voice:", openai_voice)
        else:
            print("Using Google")

        # Initialize COM library
        comtypes.CoInitialize()
        # Convert PPT slides to images using COM interface (Windows only)
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        pp_constants = comtypes.client.Constants(powerpoint)
        powerpoint.Visible = 1
        ppt = powerpoint.Presentations.Open(os.path.abspath(ppt_path))

        # Export slides as images
        # ppt.SaveAs(os.path.abspath(TEMP_IMAGES_FOLDER), pp_constants.ppSaveAsPNG)  # 17 corresponds to JPG format
        ppt.SaveAs(os.path.abspath(TEMP_IMAGES_FOLDER), pp_constants.ppSaveAsPNG )  # 18 corresponds to PNG format
        ppt.Close()
        powerpoint.Quit()

        # Extract videos from slides
        videos = extract_videos_from_slides(ppt_path)

        # Load presentation to get slide notes
        prs = Presentation(ppt_path)

        for idx, slide in enumerate(prs.slides):
            # Skip slides not in the list
            if idx+1 not in slides_numbers_to_process and len(slides_numbers_to_process) > 0:
                continue

            # Get slide notes
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            slide_image_path = f"{TEMP_IMAGES_FOLDER}/slide{idx+1}.PNG"
            duration = min_time_per_slide if min_time_per_slide is not None else 1
            slide_ai_voice = openai_voice # by default, unless overwritten

            # get cleaned notes, and all settings from slide
            slide_settings, notes = get_slide_settings(remove_comments(notes))
            min_time_from_notes = slide_settings.get("min_time", None)
            pause_time_at_end_from_notes = slide_settings.get("pause_time_at_end", None)
            ai_voice_from_notes = slide_settings.get("ai_voice", None)

            # overwrite if needed
            min_time_per_slide = min_time_from_notes if min_time_from_notes is not None else min_time_per_slide
            pause_time_at_end = pause_time_at_end_from_notes if pause_time_at_end_from_notes is not None else pause_time_at_end
            slide_ai_voice = ai_voice_from_notes if ai_voice_from_notes is not None else openai_voice

            print("Slide", idx+1, "Settings :", slide_settings, "Notes:", notes)

            if slide_ai_voice not in  ["alloy", "echo", "fable", "onyx", "nova", "shimmer"]:
                print("Invalid voice selected. Using default voice:", openai_voice)
                slide_ai_voice = openai_voice

            # Convert notes to speech
            if notes is not None and notes.strip():
                audio_path = f"{TEMP_AUDIO_FOLDER}/audio_{idx+1}.mp3"

                if provider == "openai":
                    # Generate speech using OpenAI
                    print("Generating speech using OpenAI with voice:", slide_ai_voice)
                    response = openai_client.audio.speech.create(
                        model="tts-1-hd",  # or "tts-1-hd" for higher quality
                        voice=slide_ai_voice,  # options: alloy, echo, fable, onyx, nova, shimmer
                        input=notes,
                        response_format="mp3"
                    )
                    
                    # Save the audio file
                    response.stream_to_file(audio_path)
                else: # default / google
                    tts = gTTS(text=notes, lang=language, tld=accent)
                    tts.save(audio_path)

                # if speed_factor != 1:
                #     print("Speeding up audio by", speed_factor)
                #     output_path = f"{TEMP_AUDIO_FOLDER}/audio_spedup_{idx+1}.mp3"
                #     speed_up_audio_ffmpeg(input_path=audio_path, output_path=output_path, speed_factor=speed_factor)
                #     audio_path = output_path # overwrite audio path

                audio_clip = AudioFileClip(audio_path)
                duration = audio_clip.duration
            else:
                audio_clip = None
                duration = 1  # Default duration if no notes

            # Ensure minimum time per slide
            if min_time_from_notes is not None:
                duration = max(duration, min_time_from_notes)
            elif min_time_per_slide > 0:
                duration = max(duration, min_time_per_slide)

            # Add pause at end if needed
            if pause_time_at_end > 0:
                duration += pause_time_at_end

            # Check for embedded videos
            detected_video = None
            for shape in slide.shapes:
            # breakpoint()
            # Media shapes may have .media or .src attributes
                if shape.shape_type == MSO_SHAPE_TYPE.MEDIA and shape.media_type == PP_MEDIA_TYPE.MOVIE:
                    # Get the relationship ID of the media
                    # rId = shape._element.xpath('.//a:videoFile')[0].attrib
                
                    # iterate over the slide's relationships
                    num_rels = 5
                    for i in range(1, num_rels+1):
                        curr_index = f"rId{i}"
                        try:
                            filename = os.path.basename(slide.part.rels[curr_index].target_partname)
                            
                            if any(filename in path for path in videos):
                                print(f"Found video: {slide.part.rels[curr_index].target_partname} for slide: {idx + 1}")
                                detected_video = filename
                                break
                                
                        except Exception as e:
                                pass
                                # print(f"Error while scanning for videos for slide {slide_index}: {e}")

            # right now we don't support video + audio from notes. so either or
            if detected_video is not None:
                clip = VideoFileClip(os.path.join(TEMP_VIDEOS_FOLDER, detected_video))
                clip.duration = max(duration, clip.duration)

                if speed_factor > 1:
                    clip.duration = clip.duration * speed_factor

            else:
                # Create video clip from image
                clip = ImageClip(slide_image_path).set_duration(duration)
                if audio_clip:
                    clip = clip.set_audio(audio_clip)

                if speed_factor > 1:
                    clip.duration = clip.duration * speed_factor
            
            if clip.audio is not None:
                clip = clip.set_audio(clip.audio.set_fps(44100))
            clips.append(clip)

        # Concatenate all clips and write the final video
        if clips:
            # introduce a little crossfade:
            crossfade_duration = 0.3
            for i in range(1, len(clips)):
                clips[i] = clips[i].crossfadein(crossfade_duration)

            final_clip = concatenate_videoclips(clips, method="compose")
            final_clip = final_clip.set_fps(fps)
            if hasattr(final_clip, 'audio') and final_clip.audio is not None:
                final_clip.audio = final_clip.audio.set_fps(44100)  # Standard audio sample rate
            
            temp_name = "_temp.mp4"
            final_clip.write_videofile(output_dir + "/" + temp_name, fps=fps, audio_codec="aac", audio_bitrate='192k', codec="libx264", preset="ultrafast", threads=4)
            final_clip.close()

            # speed up video
            if speed_factor != 1:
                print("Speeding up video by", speed_factor)
                speed_up_video_ffmpeg(input_path=output_dir + "/" + temp_name, output_path=output_dir + "/" + output_video, speed_factor=speed_factor)
                # delete the temp file
                os.remove(output_dir + "/" + temp_name)
            else:
                os.rename(output_dir + "/" + temp_name, output_dir + "/" + output_video)
                
    except Exception as e:
        print(f"Error: {e}")
    finally:
        # Clean up resources
        for clip in clips:
            clip.close()
        cleanup_temp_dirs()   

if __name__ == "__main__":
    # if len(sys.argv) > 1:
    #     ppt_path = sys.argv[1]
    #     convert_ppt_to_video(ppt_path)
    # else:
    #     print("Usage: python ppt_to_video.py <path_to_ppt_file>")

    convert_ppt_to_video("testppt_video.pptx")